"""
generate_vplmn_ppt.py

Reads VPLMN/MVNO outage data from an Excel file and produces a PowerPoint
(.pptx) report with one slide per Error Type (e.g. "Lost Service", "System
failure"), containing one native, editable bar chart per Roaming Partner
(MVNO) on that slide.

Expected input columns (case-sensitive):
    - "VPLMN"       : the visited network that went down (free text)
    - "MVNO"        : one or more roaming partners, e.g. "BICS" or
                       "BICS, JT" or "BICS/JT" (comma/slash/ampersand separated)
    - "Error Type"  : the failure category, e.g. "Lost Service" or
                       "System failure". Rows that mention more than one
                       category (e.g. "System failure, Lost Service") are
                       assigned to whichever category is mentioned FIRST,
                       so a single incident is never double-counted across
                       categories.
    - "Duration"    : free-text duration in almost any common format:
                       "12 Minutes", "1 hr 10 mins", "01 hrs 50 Minutes",
                       "2 Hours 5 Minutes", "10 mins", "45", "2+4 Minutes", ...

Usage:
    python generate_vplmn_ppt.py
    python generate_vplmn_ppt.py --input data/VPLMN_MVNO_table.xlsx --output output/VPLMN_Outage_Report.pptx --top-n 10

Requires: pandas, python-pptx, openpyxl
    pip install pandas python-pptx openpyxl
"""

import argparse
import re
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --------------------------------------------------------------------------
# Configuration (overridable via command-line flags, see main() below)
# --------------------------------------------------------------------------
DEFAULT_INPUT_PATH = "data/VPLMN_MVNO_table.xlsx"
DEFAULT_OUTPUT_PATH = "output/VPLMN_Outage_Report.pptx"
DEFAULT_TOP_N = 10          # max VPLMNs shown per chart (None = show all)
CHARTS_PER_ROW = 3          # layout: how many MVNO charts per row on a slide

TITLE_COLOR = RGBColor(0x00, 0x57, 0xAD)
LABEL_COLOR = RGBColor(0x2E, 0x75, 0xB6)
BAR_COLOR = RGBColor(0x4B, 0xAC, 0xC6)
SUBTITLE_COLOR = RGBColor(0x6B, 0x72, 0x80)
CHART_BORDER_COLOR = RGBColor(0x40, 0x40, 0x40)

# Optional descriptive subheading shown under the slide title for each error
# type. Add an entry here for any Error Type label you expect to see; types
# not listed fall back to a generic auto-generated subheading.
SUBTITLES = {
    "Lost Service": "Sudden service disruption due to Roaming Partner / MNO outage, or customer-side power/network outage.",
    "System failure": "Delayed Registration of Devices in Network results in System Failure Errors",
}

# Known error-type keywords to detect, in the order they should be checked.
# A row is assigned to whichever of these keywords appears EARLIEST in its
# "Error Type" text, so multi-label rows (e.g. "System failure, Lost
# Service") are counted exactly once, never in both buckets.
ERROR_TYPE_KEYWORDS = {
    "Lost Service": ["lost service", "lost-service"],
    "System failure": ["system failure"],
}

# Optional manual aliases for MVNO names that refer to the same roaming
# partner but are spelled differently in the source data (beyond simple
# case differences, which are handled automatically). Edit this for your
# own dataset as needed -- keys are matched case-insensitively.
MVNO_ALIASES = {
    "orange": "Orange_901",
}

# Corporate-entity suffixes that sometimes get separated from the company
# name by a comma (e.g. "T-Mobile USA, Inc" or "A1 Slovenija, d.d."). When
# a VPLMN cell is split on commas (see split_vplmn_list below), a fragment
# matching one of these is re-attached to the preceding name instead of
# being treated as its own (bogus) VPLMN entry.
CORPORATE_SUFFIXES = {"inc", "inc.", "s.a.", "s.a.u.", "sa", "sau", "d.d.", "dd", "ltd", "ltd."}

# KNOWN LIMITATION: a small number of source rows list several affected
# VPLMNs in one cell separated by SPACES ONLY, with no comma at all, e.g.:
#   "AT&T Mobility T-Mobile USA, Inc Orange France Bouygues Telecom ..."
# There is no reliable, generic way to split that back into individual
# network names. If you hit rows like this, add the exact cell text here
# mapped to the correct list of names, e.g.:
#   KNOWN_VPLMN_SPLITS = {
#       "AT&T Mobility T-Mobile USA, Inc Orange France ...": [
#           "AT&T Mobility", "T-Mobile USA, Inc", "Orange France", ...
#       ],
#   }
# Rows not listed here that can't be comma-split will be kept as a single
# (long, clearly wrong-looking) VPLMN entry rather than silently guessed at.
KNOWN_VPLMN_SPLITS: dict = {}


# --------------------------------------------------------------------------
# Text encoding fix-up
# --------------------------------------------------------------------------
def fix_mojibake(value):
    """
    Repair text that was double-encoded (UTF-8 bytes misread as
    cp1252/latin-1), which shows up as e.g. "TelefÃ³nica" instead of
    "Telefónica". Leaves the value unchanged if it isn't affected.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return value
    text = str(value)
    try:
        return text.encode("cp1252").decode("utf-8")
    except (UnicodeEncodeError, UnicodeDecodeError):
        return text


# --------------------------------------------------------------------------
# Duration parsing
# --------------------------------------------------------------------------
def parse_duration_to_minutes(value) -> float:
    """
    Convert a free-text duration into total minutes.

    Handles (case-insensitively):
        "45 Minutes", "45 mins", "45 min", "45"
        "1 hr 10 mins", "1 hour 10 minutes", "01 hrs 50 Minutes"
        "2+4 Minutes"  -> treated as 2 + 4 = 6 minutes
        combinations/spacing/pluralization variants of the above

    Returns 0.0 for missing/unparseable values.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return 0.0

    text = str(value).strip().lower()
    if not text:
        return 0.0

    # "2+4 Minutes" style: sum every number found, no hour component assumed.
    if "+" in text:
        numbers = re.findall(r"\d+", text)
        return float(sum(int(n) for n in numbers))

    hours = sum(int(n) for n in re.findall(r"(\d+)\s*(?:hours?|hrs?|hr)\b", text))
    minutes = sum(int(n) for n in re.findall(r"(\d+)\s*(?:minutes?|mins?|min)\b", text))

    if hours == 0 and minutes == 0:
        # Fallback: a bare number with no unit at all is assumed to already
        # be in minutes (e.g. a cell that just contains "45").
        bare_number = re.match(r"^(\d+(?:\.\d+)?)$", text)
        if bare_number:
            return float(bare_number.group(1))
        return 0.0

    return float(hours * 60 + minutes)


# --------------------------------------------------------------------------
# MVNO / Error Type / VPLMN normalization
# --------------------------------------------------------------------------
def split_mvnos(value, canon_map: dict) -> list:
    """
    Split a (possibly multi-valued) MVNO cell into a clean list of names,
    normalizing case so "BICS", "bics" and "Bics" are treated as the same
    partner, and applying MVNO_ALIASES for known alternate spellings.
    `canon_map` accumulates {uppercase_key: display_name} the first time
    each partner is seen, so the display casing stays consistent report-wide.
    Handles separators: comma, slash, ampersand, and the word "and".
    A single incident naming several MVNOs is attributed to each of them.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return []
    text = str(value)
    text = re.sub(r"\(.*?\)", " ", text)          # drop parenthetical asides
    text = re.sub(r"[&/]|(?<!\w)and(?!\w)", ",", text, flags=re.IGNORECASE)
    parts = [p.strip() for p in text.split(",") if p.strip()]

    result = []
    for part in parts:
        key = part.strip().lower()
        display = MVNO_ALIASES.get(key, part.strip())
        canon_key = display.upper()
        if canon_key not in canon_map:
            canon_map[canon_key] = display
        result.append(canon_map[canon_key])
    return result


def classify_error_type(value):
    """
    Return exactly one error-type label for a row (or None if the text
    doesn't match any known keyword). If multiple keywords are present,
    the one appearing earliest in the text wins -- this prevents a single
    incident from being counted under more than one error type.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return None
    text = str(value).lower()

    best_label, best_pos = None, None
    for label, keywords in ERROR_TYPE_KEYWORDS.items():
        for kw in keywords:
            pos = text.find(kw)
            if pos != -1 and (best_pos is None or pos < best_pos):
                best_label, best_pos = label, pos
    return best_label


def split_vplmn_list(value) -> list:
    """
    Clean a VPLMN cell and split it into one or more individual network
    names. Some source rows list several affected VPLMNs in one cell,
    comma-separated (e.g. "ROAMING PARTNER DOWN" style alerts naming every
    affected network) -- this splits those into separate entries.

    Naively splitting on every comma would also break apart names that
    legitimately contain one, such as "T-Mobile USA, Inc" or
    "A1 Slovenija, d.d." -- fragments matching CORPORATE_SUFFIXES are
    re-attached to the preceding name instead of becoming their own entry.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return []
    text = re.sub(r"\[.*?\]", "", str(value))  # strip bracketed annotations
    text = text.strip()

    if text in KNOWN_VPLMN_SPLITS:
        return list(KNOWN_VPLMN_SPLITS[text])

    raw_parts = [p.strip() for p in text.split(",")]

    merged = []
    for part in raw_parts:
        if part.lower() in CORPORATE_SUFFIXES:
            if merged:
                merged[-1] = f"{merged[-1]}, {part}"
            continue
        if part:
            merged.append(part)
    return merged


# --------------------------------------------------------------------------
# Data pipeline
# --------------------------------------------------------------------------
def load_and_aggregate(input_path: Path) -> pd.DataFrame:
    """
    Read the raw Excel file and return a tidy, aggregated DataFrame with
    columns: ErrorType, MVNO, VPLMN, Incidents, TotalDurationMin.
    """
    df = pd.read_excel(input_path)

    required_cols = {"VPLMN", "MVNO", "Error Type", "Duration"}
    missing = required_cols - set(df.columns)
    if missing:
        raise ValueError(f"Input file is missing required column(s): {sorted(missing)}")

    df["Error Type"] = df["Error Type"].apply(fix_mojibake)
    df["MVNO"] = df["MVNO"].apply(fix_mojibake)
    df["VPLMN"] = df["VPLMN"].apply(fix_mojibake)

    mvno_canon_map: dict = {}
    df["ErrorType"] = df["Error Type"].apply(classify_error_type)
    df["MVNOs"] = df["MVNO"].apply(lambda v: split_mvnos(v, mvno_canon_map))
    df["VPLMNs"] = df["VPLMN"].apply(split_vplmn_list)
    df["DurationMin"] = df["Duration"].apply(parse_duration_to_minutes)

    exploded_rows = []
    for _, row in df.iterrows():
        if row["ErrorType"] is None or not row["VPLMNs"] or not row["MVNOs"]:
            continue
        for mvno in row["MVNOs"]:
            for vplmn in row["VPLMNs"]:
                exploded_rows.append(
                    {
                        "ErrorType": row["ErrorType"],
                        "MVNO": mvno,
                        "VPLMN": vplmn,
                        "DurationMin": row["DurationMin"],
                    }
                )

    if not exploded_rows:
        raise ValueError("No usable rows found after cleaning -- check the input data.")

    exploded = pd.DataFrame(exploded_rows)
    agg = (
        exploded.groupby(["ErrorType", "MVNO", "VPLMN"])
        .agg(Incidents=("DurationMin", "count"), TotalDurationMin=("DurationMin", "sum"))
        .reset_index()
    )
    return agg


# --------------------------------------------------------------------------
# PowerPoint report generation
# --------------------------------------------------------------------------
def add_chart_area_shadow(chart):
    """
    Apply a soft drop-shadow effect to the chart area itself (not the
    surrounding border box). Charts are graphic frames, not plain shapes,
    so python-pptx has no high-level shadow API for them -- this adds a
    <c:spPr><a:effectLst><a:outerShdw> element directly to the chart's own
    XML (as a sibling of <c:chart>, per the chartSpace schema).
    """
    c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    chart_space = chart._chartSpace

    c_spPr = chart_space.makeelement(f"{{{c_ns}}}spPr", {})
    a_effectLst = chart_space.makeelement(f"{{{a_ns}}}effectLst", {})
    a_outerShdw = chart_space.makeelement(
        f"{{{a_ns}}}outerShdw",
        {"blurRad": "50800", "dist": "25400", "dir": "5400000", "rotWithShape": "0"},
    )
    a_srgbClr = chart_space.makeelement(f"{{{a_ns}}}srgbClr", {"val": "000000"})
    a_alpha = chart_space.makeelement(f"{{{a_ns}}}alpha", {"val": "35000"})

    a_srgbClr.append(a_alpha)
    a_outerShdw.append(a_srgbClr)
    a_effectLst.append(a_outerShdw)
    c_spPr.append(a_effectLst)

    # <c:spPr> must come immediately after <c:chart> in the chartSpace schema
    c_chart_el = chart_space.find(f"{{{c_ns}}}chart")
    c_chart_el.addnext(c_spPr)


def add_chart_to_slide(slide, left, top, width, height, mvno_label, cats, vals):
    """Add one labeled, bordered, native bar chart to the slide at the given position."""
    label_h = Inches(0.35)

    title_box = slide.shapes.add_textbox(left, top, width, label_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = mvno_label
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = LABEL_COLOR
    p.font.name = "Arial"

    chart_top = top + label_h
    chart_height = height - label_h

    chart_data = CategoryChartData()
    # reverse so the largest value is drawn at the top of the horizontal bar chart
    chart_data.categories = list(reversed(cats))
    chart_data.add_series("Duration (min)", list(reversed(vals)))

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        left, chart_top, width, chart_height,
        chart_data,
    )
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = "Duration (min)"

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = "#,##0"
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(10)

    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BAR_COLOR

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "Minutes"

    # Draw a thin border rectangle around the chart's plot box, on top of
    # the chart so the outline stays crisp, with no fill so the chart shows
    # through underneath.
    border_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, chart_top, width, chart_height)
    # Subtle corner rounding (default rounded-rect preset is quite round; tone it down)
    try:
        border_box.adjustments[0] = 0.045
    except (IndexError, AttributeError):
        pass
    border_box.fill.background()
    border_box.line.color.rgb = CHART_BORDER_COLOR
    border_box.line.width = Pt(1)
    border_box.shadow.inherit = False

    add_chart_area_shadow(chart)


def add_error_type_slide(prs, agg, error_type, top_n):
    """Add one slide for a given error type, with one bar chart per MVNO."""
    subset = agg[agg["ErrorType"] == error_type]
    partners = sorted(subset["MVNO"].unique())
    if not partners:
        return

    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    slide_w, slide_h = prs.slide_width, prs.slide_height

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), slide_w - Inches(0.8), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = f"{error_type} (MNO Perspective)"
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = TITLE_COLOR
    tp.font.name = "Arial"

    subtitle_text = SUBTITLES.get(error_type, f"{error_type} incidents by Roaming Partner (MVNO)")
    subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.82), slide_w - Inches(0.8), Inches(0.4))
    sp = subtitle_box.text_frame.paragraphs[0]
    sp.text = subtitle_text
    sp.font.size = Pt(13)
    sp.font.color.rgb = SUBTITLE_COLOR
    sp.font.name = "Arial"

    n = len(partners)
    cols = min(CHARTS_PER_ROW, n)
    rows = -(-n // cols)  # ceil division

    margin = Inches(0.4)
    gap = Inches(0.25)
    top0 = Inches(1.45)
    grid_w = slide_w - 2 * margin
    grid_h = slide_h - top0 - Inches(0.3)
    cell_w = (grid_w - gap * (cols - 1)) // cols
    cell_h = (grid_h - gap * (rows - 1)) // rows

    for idx, mvno in enumerate(partners):
        r, c = divmod(idx, cols)
        left = margin + c * (cell_w + gap)
        top = top0 + r * (cell_h + gap)

        partner_data = subset[subset["MVNO"] == mvno].sort_values("TotalDurationMin", ascending=False)
        total_vplmns = len(partner_data)
        if top_n is not None and total_vplmns > top_n:
            partner_data = partner_data.head(top_n)
            label = f"{mvno}  (top {top_n} of {total_vplmns} VPLMN)"
        else:
            label = mvno

        cats = partner_data["VPLMN"].tolist()
        vals = partner_data["TotalDurationMin"].tolist()
        add_chart_to_slide(slide, left, top, cell_w, cell_h, label, cats, vals)


def generate_pptx_report(agg: pd.DataFrame, output_path: Path, top_n):
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # widescreen 16:9
    prs.slide_height = Inches(7.5)

    error_types = sorted(agg["ErrorType"].unique())
    for error_type in error_types:
        add_error_type_slide(prs, agg, error_type, top_n)

    prs.save(output_path)
    print(f"Saved report: {output_path}  ({len(error_types)} slide(s))")


# --------------------------------------------------------------------------
# Entry point
# --------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(description="Generate a VPLMN/MVNO outage PowerPoint report.")
    parser.add_argument("--input", default=DEFAULT_INPUT_PATH, help="Path to the source .xlsx file")
    parser.add_argument("--output", default=DEFAULT_OUTPUT_PATH, help="Path to write the .pptx report")
    parser.add_argument(
        "--top-n",
        type=int,
        default=DEFAULT_TOP_N,
        help="Max VPLMNs to show per chart (omit/-1 for no limit)",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output)
    top_n = None if args.top_n is None or args.top_n < 0 else args.top_n

    if not input_path.exists():
        sys.exit(f"Input file not found: {input_path}")

    agg = load_and_aggregate(input_path)
    generate_pptx_report(agg, output_path, top_n)


if __name__ == "__main__":
    main()